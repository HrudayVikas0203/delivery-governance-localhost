from io import BytesIO
from pathlib import Path
from types import SimpleNamespace
import uuid

from fastapi.testclient import TestClient
from pptx import Presentation

from app.db.session import SessionLocal
from app.main import app
from app.models.status import GeneratedReport, ReportFormat, ReportTemplate, ReportType


PPTX_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _headers(client: TestClient) -> dict[str, str]:
    response = client.post("/api/v1/auth/login", json={"email": "praveen.baburaya@delta.com", "password": "Demo@123"})
    assert response.status_code == 200, response.text
    return {"Authorization": f"Bearer {response.json()['access_token']}"}


def _account(client: TestClient, headers: dict[str, str]) -> str:
    employees = client.get("/api/v1/governance/employees", headers=headers).json()
    program_manager_id = next(employee["id"] for employee in employees if employee["role"] == "program_manager")
    response = client.post(
        "/api/v1/governance/accounts",
        headers=headers,
        json={"name": f"Template Storage {uuid.uuid4().hex}", "industry": "Technology", "country": "USA", "business_unit": "Delivery", "program_manager_id": program_manager_id},
    )
    assert response.status_code == 201, response.text
    return response.json()["id"]


def _pptx(marker: str = "{{PROJECT_NAME}}") -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(1, 1, 8, 1).text = marker
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_valid_template_is_persisted_and_reopens_after_new_session():
    client = TestClient(app)
    headers = _headers(client)
    account_id = _account(client, headers)
    content = _pptx("PERSISTED_TEMPLATE")

    response = client.post(
        f"/api/v1/governance/accounts/{account_id}/template",
        headers=headers,
        files={"file": ("account.pptx", content, PPTX_TYPE)},
    )
    assert response.status_code == 201, response.text
    template_id = response.json()["id"]

    with SessionLocal() as first_session:
        first = first_session.get(ReportTemplate, template_id)
        assert first is not None
        assert first.file_path == f"database://report-templates/{template_id}"
        assert bytes(first.content_bytes) == content
    with SessionLocal() as restarted_session:
        stored = restarted_session.get(ReportTemplate, template_id)
        assert stored is not None
        reopened = Presentation(BytesIO(bytes(stored.content_bytes)))
        assert len(reopened.slides) == 1


def test_invalid_extension_and_corrupt_pptx_are_rejected_without_replacing_existing():
    client = TestClient(app)
    headers = _headers(client)
    account_id = _account(client, headers)
    valid = client.post(
        f"/api/v1/governance/accounts/{account_id}/template",
        headers=headers,
        files={"file": ("valid.pptx", _pptx("ORIGINAL"), PPTX_TYPE)},
    )
    assert valid.status_code == 201, valid.text
    template_id = valid.json()["id"]

    wrong_extension = client.post(
        f"/api/v1/governance/accounts/{account_id}/template",
        headers=headers,
        files={"file": ("invalid.pdf", b"not a presentation", "application/pdf")},
    )
    assert wrong_extension.status_code == 400
    assert ".pptx" in wrong_extension.json()["detail"]

    corrupt = client.post(
        f"/api/v1/governance/accounts/{account_id}/template",
        headers=headers,
        files={"file": ("corrupt.pptx", b"not a presentation", PPTX_TYPE)},
    )
    assert corrupt.status_code == 400
    assert "invalid or corrupted" in corrupt.json()["detail"]

    with SessionLocal() as db:
        stored = db.get(ReportTemplate, template_id)
        assert stored is not None
        assert stored.filename == "valid.pptx"
        assert b"not a presentation" != bytes(stored.content_bytes)


def test_oversized_template_is_rejected(monkeypatch):
    settings = SimpleNamespace(ppt_template_max_bytes=100)
    monkeypatch.setattr("app.api.v1.governance.get_settings", lambda: settings)
    monkeypatch.setattr("app.services.template_storage.get_settings", lambda: settings)
    client = TestClient(app)
    headers = _headers(client)
    account_id = _account(client, headers)

    response = client.post(
        f"/api/v1/governance/accounts/{account_id}/template",
        headers=headers,
        files={"file": ("large.pptx", _pptx(), PPTX_TYPE)},
    )
    assert response.status_code == 400
    assert "size limit" in response.json()["detail"]


def test_account_creation_with_invalid_template_is_atomic():
    client = TestClient(app)
    headers = _headers(client)
    account_name = f"Atomic Account {uuid.uuid4().hex}"
    employees = client.get("/api/v1/governance/employees", headers=headers).json()
    program_manager_id = next(employee["id"] for employee in employees if employee["role"] == "program_manager")
    response = client.post(
        "/api/v1/governance/accounts/with-template",
        headers=headers,
        data={"account_data": '{"name":"' + account_name + '","industry":"Technology","country":"USA","business_unit":"Delivery","program_manager_id":"' + program_manager_id + '"}'},
        files={"file": ("invalid.pptx", b"invalid", PPTX_TYPE)},
    )
    assert response.status_code == 400
    accounts = client.get("/api/v1/governance/accounts", headers=headers).json()
    assert all(account["name"] != account_name for account in accounts)


def test_download_uses_persisted_report_bytes_when_local_output_is_missing():
    client = TestClient(app)
    headers = _headers(client)
    content = _pptx("DOWNLOAD")
    with SessionLocal() as db:
        report = GeneratedReport(
            title="Persisted Download",
            report_type=ReportType.PROJECT_REPORT,
            report_format=ReportFormat.PPTX,
            scope="test",
            status="ready",
            filename="persisted.pptx",
            content_type=PPTX_TYPE,
            size_bytes=len(content),
            content_bytes=content,
            file_path=str(Path("missing") / "persisted.pptx"),
        )
        db.add(report)
        db.commit()
        report_id = report.id

    response = client.get(f"/api/v1/reports/{report_id}/download", headers=headers)
    assert response.status_code == 200
    assert response.headers["content-type"] == PPTX_TYPE
    assert len(Presentation(BytesIO(response.content)).slides) == 1
