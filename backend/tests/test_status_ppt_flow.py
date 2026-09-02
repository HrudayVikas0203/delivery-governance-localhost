import io
import os
import tempfile
import uuid
from pathlib import Path
from types import SimpleNamespace

os.environ.setdefault("DATABASE_BACKEND", "sqlite")

from fastapi.testclient import TestClient
from pptx import Presentation

from app.core.config import get_settings
from app.db.seed import seed
from app.db.schema import ensure_schema_upgrades
from app.db.session import Base, SessionLocal, engine
from app.models.delivery import Account, Project
from app.models.status import WeeklyStatus, ReportTemplate
from app.models.tasks import Task
from app.main import app


def _management_ids(client: TestClient, headers: dict[str, str]) -> tuple[str, str, str]:
    employees = client.get("/api/v1/governance/employees", headers=headers).json()
    program_manager = next(row["id"] for row in employees if row["role"] == "program_manager")
    project_manager = next(row["id"] for row in employees if row["role"] == "project_manager" and row["manager_id"] == program_manager)
    studio_head = next(row["id"] for row in employees if row["role"] in {"delivery_head", "studio_head"})
    return program_manager, project_manager, studio_head


def test_account_project_allocation_status_and_ppt_flow() -> None:
    get_settings.cache_clear()
    Base.metadata.create_all(bind=engine)
    ensure_schema_upgrades()
    seed()
    suffix = uuid.uuid4().hex[:8]
    account_name = f"Northwind Capital {suffix}"
    project_name = f"Treasury Analytics {suffix}"

    client = TestClient(app)
    resp = client.post(
        "/api/v1/auth/login",
        json={"email": "praveen.baburaya@delta.com", "password": "Demo@123"},
    )
    assert resp.status_code == 200, resp.text
    token = resp.json()["access_token"]
    headers = {"Authorization": f"Bearer {token}"}
    program_manager_id, project_manager_id, studio_head_id = _management_ids(client, headers)

    account_resp = client.post(
        "/api/v1/governance/accounts",
        json={
            "name": account_name,
            "industry": "Finance",
            "country": "USA",
            "business_unit": "Private Banking",
            "contract_value": 1500000,
            "start_date": "2026-08-01",
            "end_date": "2027-08-01",
            "program_manager_id": program_manager_id,
        },
        headers=headers,
    )
    assert account_resp.status_code == 201, account_resp.text
    account_id = account_resp.json()["id"]

    project_resp = client.post(
        "/api/v1/governance/projects",
        json={
            "account_id": account_id,
            "name": project_name,
            "phase": "development",
            "client": account_name,
            "budget_used": 0,
            "budget_total": 1200000,
            "tech_stack": ["React", "FastAPI", "Postgres"],
            "sprint_number": 3,
            "description": "Modern treasury analytics and risk dashboard.",
            "start_date": "2026-08-05",
            "completion_percent": 55,
            "project_manager_id": project_manager_id,
        },
        headers=headers,
    )
    assert project_resp.status_code == 201, project_resp.text
    project_id = project_resp.json()["id"]

    employee_resp = client.get("/api/v1/governance/employees", headers=headers)
    assert employee_resp.status_code == 200, employee_resp.text
    employee_id = next(e["id"] for e in employee_resp.json() if e["email"] == "deepak.sharma@delta.com")

    allocate_resp = client.post(
        "/api/v1/governance/allocations",
        json={
            "project_id": project_id,
            "employee_id": studio_head_id,
            "allocation_role": "developer",
            "allocation_percent": 80,
            "start_date": "2026-08-05",
            "reporting_manager_id": None,
        },
        headers=headers,
    )
    assert allocate_resp.status_code == 201, allocate_resp.text

    status_resp = client.post(
        "/api/v1/governance/status",
        json={
            "employee_id": studio_head_id,
            "project_id": project_id,
            "week_start": "2026-08-10",
            "status": "submitted",
            "fields": {
                "project": project_name,
                "account": account_name,
                "reportingFrequency": "Weekly",
                "overallStatus": "Green",
                "completionPercent": 61,
                "hoursWorked": 42,
                "achievements": "Delivered the treasury dashboard wireframes and integrated core risk calculations.",
                "blockers": "No blockers. Dependency on upstream market data feed remains under watch.",
                "risks": "Market data source latency may affect daily refresh windows.",
                "nextWeekPlan": "Complete reconciliation workflow and start user acceptance testing.",
                "supportRequired": "Need confirmation on final data refresh SLA from the client team.",
            },
        },
        headers=headers,
    )
    assert status_resp.status_code == 201, status_resp.text

    template_path = Path(tempfile.gettempdir()) / f"flow-template-{uuid.uuid4().hex}.pptx"
    template = Presentation()
    template_slide = template.slides.add_slide(template.slide_layouts[6])
    template_slide.shapes.add_textbox(1.0, 1.0, 8.0, 1.0).text = "{{PROJECT_NAME}}"
    template.save(template_path)
    with template_path.open("rb") as file_handle:
        template_resp = client.post(
            f"/api/v1/governance/accounts/{account_id}/template",
            files={"file": (template_path.name, file_handle, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            headers=headers,
        )
    assert template_resp.status_code == 201, template_resp.text

    report_resp = client.post(
        "/api/v1/reports",
        json={
            "title": f"{project_name} Weekly Status",
            "report_type": "project_report",
            "report_format": "pptx",
            "scope": f"project:{project_id}",
            "account_id": account_id,
            "project_id": project_id,
            "status_frequency": "weekly",
            "use_celery": False,
            "llm": None,
        },
        headers=headers,
    )
    assert report_resp.status_code == 201, report_resp.text
    report = report_resp.json()
    assert report["status"] == "ready", report
    assert report["file_path"], report

    file_path = Path(report["file_path"])
    assert file_path.exists(), file_path
    assert file_path.suffix == ".pptx"


def test_account_template_is_auto_selected_for_project_reports() -> None:
    get_settings.cache_clear()
    Base.metadata.create_all(bind=engine)
    ensure_schema_upgrades()
    seed()

    client = TestClient(app)
    login_resp = client.post(
        "/api/v1/auth/login",
        json={"email": "praveen.baburaya@delta.com", "password": "Demo@123"},
    )
    assert login_resp.status_code == 200, login_resp.text
    token = login_resp.json()["access_token"]
    headers = {"Authorization": f"Bearer {token}"}
    program_manager_id, project_manager_id, _ = _management_ids(client, headers)

    account_name = f"Northwind Capital {uuid.uuid4().hex[:8]}"
    account_resp = client.post(
        "/api/v1/governance/accounts",
        json={
            "name": account_name,
            "industry": "Finance",
            "country": "USA",
            "business_unit": "Private Banking",
            "contract_value": 1500000,
            "start_date": "2026-08-01",
            "end_date": "2027-08-01",
            "program_manager_id": program_manager_id,
        },
        headers=headers,
    )
    assert account_resp.status_code == 201, account_resp.text
    account_id = account_resp.json()["id"]

    project_name = f"Treasury Analytics {uuid.uuid4().hex[:8]}"
    project_resp = client.post(
        "/api/v1/governance/projects",
        json={
            "account_id": account_id,
            "name": project_name,
            "phase": "development",
            "client": account_name,
            "budget_used": 0,
            "budget_total": 1200000,
            "tech_stack": ["React", "FastAPI", "Postgres"],
            "sprint_number": 3,
            "description": "Modern treasury analytics and risk dashboard.",
            "start_date": "2026-08-05",
            "completion_percent": 55,
            "project_manager_id": project_manager_id,
        },
        headers=headers,
    )
    assert project_resp.status_code == 201, project_resp.text
    project_id = project_resp.json()["id"]

    template_path = Path(tempfile.gettempdir()) / f"{uuid.uuid4().hex}.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "REPORT TITLE"
    body = slide.shapes.add_textbox(1.5, 2.0, 8, 2.2)
    body.text_frame.text = "EXECUTIVE SUMMARY"
    prs.save(template_path)

    with template_path.open("rb") as file_handle:
        template_resp = client.post(
            f"/api/v1/governance/accounts/{account_id}/template",
            files={"file": (template_path.name, file_handle, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            headers=headers,
        )
    assert template_resp.status_code == 201, template_resp.text
    template_id = template_resp.json()["id"]

    report_resp = client.post(
        "/api/v1/reports",
        json={
            "title": f"{project_name} Weekly Status",
            "report_type": "project_report",
            "report_format": "pptx",
            "scope": f"project:{project_id}",
            "account_id": account_id,
            "project_id": project_id,
            "status_frequency": "weekly",
            "use_celery": False,
            "llm": None,
        },
        headers=headers,
    )
    assert report_resp.status_code == 201, report_resp.text
    report = report_resp.json()
    assert report["template_id"] == template_id, report
    assert report["file_path"], report
    assert Path(report["file_path"]).exists(), report["file_path"]

    generated = Presentation(report["file_path"])
    assert len(generated.slides) >= 1
    text_chunks = []
    for slide in generated.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_chunks.append(shape.text)
    assert any(project_name in chunk for chunk in text_chunks)


def test_account_template_replace_remove_and_project_isolation() -> None:
    get_settings.cache_clear()
    Base.metadata.create_all(bind=engine)
    ensure_schema_upgrades()
    seed()

    client = TestClient(app)
    login_resp = client.post(
        "/api/v1/auth/login",
        json={"email": "praveen.baburaya@delta.com", "password": "Demo@123"},
    )
    assert login_resp.status_code == 200, login_resp.text
    headers = {"Authorization": f"Bearer {login_resp.json()['access_token']}"}
    program_manager_id, project_manager_id, _ = _management_ids(client, headers)

    def create_account(name: str) -> str:
        response = client.post(
            "/api/v1/governance/accounts",
            json={"name": name, "industry": "Finance", "country": "USA", "business_unit": "Banking", "program_manager_id": program_manager_id},
            headers=headers,
        )
        assert response.status_code == 201, response.text
        return response.json()["id"]

    def create_project(account_id: str, name: str) -> str:
        response = client.post(
            "/api/v1/governance/projects",
            json={"account_id": account_id, "name": name, "phase": "development", "client": name, "project_manager_id": project_manager_id},
            headers=headers,
        )
        assert response.status_code == 201, response.text
        return response.json()["id"]

    def make_template(path: Path, marker: str) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "{{ACCOUNT_NAME}}"
        body = slide.shapes.add_textbox(1.5, 2.0, 8, 2.2)
        body.text_frame.text = f"TEMPLATE_MARKER_{marker}"
        prs.save(path)

    account_a = create_account(f"Template A {uuid.uuid4().hex[:8]}")
    account_b = create_account(f"Template B {uuid.uuid4().hex[:8]}")
    project_a = create_project(account_a, f"Project A {uuid.uuid4().hex[:8]}")
    project_b = create_project(account_b, f"Project B {uuid.uuid4().hex[:8]}")

    template_a = Path(tempfile.gettempdir()) / f"template-a-{uuid.uuid4().hex}.pptx"
    template_a_replacement = Path(tempfile.gettempdir()) / f"template-a-replacement-{uuid.uuid4().hex}.pptx"
    template_b = Path(tempfile.gettempdir()) / f"template-b-{uuid.uuid4().hex}.pptx"
    make_template(template_a, "A")
    make_template(template_a_replacement, "A-REPLACED")
    make_template(template_b, "B")

    def upload(account_id: str, path: Path) -> str:
        with path.open("rb") as handle:
            response = client.post(
                f"/api/v1/governance/accounts/{account_id}/template",
                files={"file": (path.name, handle, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                headers=headers,
            )
        assert response.status_code == 201, response.text
        return response.json()["id"]

    template_a_id = upload(account_a, template_a)
    assert upload(account_a, template_a_replacement) == template_a_id
    template_b_id = upload(account_b, template_b)
    assert template_a_id != template_b_id

    report_a = client.post(
        "/api/v1/reports",
        json={"title": "A report", "report_type": "project_report", "report_format": "pptx", "scope": f"project:{project_a}", "project_id": project_a, "use_celery": False},
        headers=headers,
    )
    report_b = client.post(
        "/api/v1/reports",
        json={"title": "B report", "report_type": "project_report", "report_format": "pptx", "scope": f"project:{project_b}", "project_id": project_b, "use_celery": False},
        headers=headers,
    )
    assert report_a.status_code == 201, report_a.text
    assert report_b.status_code == 201, report_b.text
    assert report_a.json()["template_id"] == template_a_id
    assert report_b.json()["template_id"] == template_b_id
    text_a = "\n".join(shape.text for slide in Presentation(report_a.json()["file_path"]).slides for shape in slide.shapes if getattr(shape, "text", ""))
    text_b = "\n".join(shape.text for slide in Presentation(report_b.json()["file_path"]).slides for shape in slide.shapes if getattr(shape, "text", ""))
    assert "TEMPLATE_MARKER_A-REPLACED" in text_a
    assert "TEMPLATE_MARKER_B" not in text_a
    assert "TEMPLATE_MARKER_B" in text_b
    assert "TEMPLATE_MARKER_A-REPLACED" not in text_b

    remove_response = client.delete(f"/api/v1/governance/accounts/{account_a}/template", headers=headers)
    assert remove_response.status_code == 204, remove_response.text

    account_response = client.get(f"/api/v1/governance/accounts/{account_a}", headers=headers)
    assert account_response.status_code == 200, account_response.text
    assert account_response.json()["ppt_template_status"] == "not_configured"

    missing_report = client.post(
        "/api/v1/reports",
        json={"title": "A missing template", "report_type": "project_report", "report_format": "pptx", "scope": f"project:{project_a}", "project_id": project_a, "use_celery": False},
        headers=headers,
    )
    assert missing_report.status_code == 400
    assert missing_report.json()["detail"] == "Account PPT template is not configured."
    assert "PackageNotFoundError" not in missing_report.text

    b_still_works = client.post(
        "/api/v1/reports",
        json={"title": "B still works", "report_type": "project_report", "report_format": "pptx", "scope": f"project:{project_b}", "project_id": project_b, "use_celery": False},
        headers=headers,
    )
    assert b_still_works.status_code == 201, b_still_works.text
    assert b_still_works.json()["template_id"] == template_b_id


def test_report_generation_retrieves_account_template_after_local_copy_is_removed(monkeypatch) -> None:
    get_settings.cache_clear()
    Base.metadata.create_all(bind=engine)
    ensure_schema_upgrades()
    seed()

    client = TestClient(app)
    login_resp = client.post(
        "/api/v1/auth/login",
        json={"email": "praveen.baburaya@delta.com", "password": "Demo@123"},
    )
    assert login_resp.status_code == 200, login_resp.text
    headers = {"Authorization": f"Bearer {login_resp.json()['access_token']}"}
    program_manager_id, project_manager_id, _ = _management_ids(client, headers)

    account_resp = client.post(
        "/api/v1/governance/accounts",
        json={"name": f"Restart Test {uuid.uuid4().hex[:8]}", "industry": "Finance", "country": "USA", "business_unit": "Banking", "program_manager_id": program_manager_id},
        headers=headers,
    )
    assert account_resp.status_code == 201, account_resp.text
    account_id = account_resp.json()["id"]

    project_resp = client.post(
        "/api/v1/governance/projects",
        json={"account_id": account_id, "name": f"Restart Project {uuid.uuid4().hex[:8]}", "phase": "development", "client": "Restart Test", "project_manager_id": project_manager_id},
        headers=headers,
    )
    assert project_resp.status_code == 201, project_resp.text
    project_id = project_resp.json()["id"]

    template_path = Path(tempfile.gettempdir()) / f"restart-template-{uuid.uuid4().hex}.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(1.0, 1.0, 8.0, 1.0).text = "{{PROJECT_NAME}}"
    prs.save(template_path)
    with template_path.open("rb") as file_handle:
        upload_resp = client.post(
            f"/api/v1/governance/accounts/{account_id}/template",
            files={"file": (template_path.name, file_handle, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            headers=headers,
        )
    assert upload_resp.status_code == 201, upload_resp.text
    template_id = upload_resp.json()["id"]

    with SessionLocal() as db:
        template = db.get(ReportTemplate, template_id)
        assert template is not None
        assert template.content_bytes
        template.file_path = str(Path(tempfile.gettempdir()) / f"missing-render-template-{uuid.uuid4().hex}.pptx")
        db.commit()
        assert not Path(template.file_path).exists()

    materialized_paths = []

    def fake_map_template(path, *_args):
        materialized_paths.append(Path(path))
        assert materialized_paths[-1].exists()
        return None, SimpleNamespace(slides=[])

    monkeypatch.setattr("app.reports.generator.map_template", fake_map_template)
    monkeypatch.setattr(
        "app.services.llm.generate_text",
        lambda *_args, **_kwargs: ('{"executive_summary":"Restart verification summary."}', "mock-model"),
    )
    report_resp = client.post(
        "/api/v1/reports",
        json={"title": "Restart report", "report_type": "project_report", "report_format": "pptx", "scope": f"project:{project_id}", "project_id": project_id, "use_celery": False, "llm": {"provider": "gemini", "model": "gemini-3.5-flash"}},
        headers=headers,
    )
    assert report_resp.status_code == 201, report_resp.text
    report = report_resp.json()
    assert report["template_id"] == template_id
    assert len(materialized_paths) == 1
    assert not materialized_paths[0].exists()
    assert Presentation(report["file_path"]).slides


def test_trimble_finance_ai_assistant_seed_data_exists() -> None:
    get_settings.cache_clear()
    Base.metadata.create_all(bind=engine)
    ensure_schema_upgrades()
    seed()

    with SessionLocal() as db:
        account = db.query(Account).filter(Account.name == "Trimble").one_or_none()
        assert account is not None, "Trimble account was not seeded"

        project = db.query(Project).filter(Project.name == "Trimble Finance AI Assistant").one_or_none()
        assert project is not None, "Trimble Finance AI Assistant project was not seeded"
        assert project.account_id == account.id

        tasks = db.query(Task).filter(Task.project_id == project.id).all()
        assert len(tasks) >= 4, "Expected multiple tasks for the Trimble project"

        weekly_statuses = db.query(WeeklyStatus).filter(WeeklyStatus.project_id == project.id).all()
        assert len(weekly_statuses) >= 4, "Expected weekly status updates for the Trimble project"

        template = db.query(ReportTemplate).filter(
            ReportTemplate.account_id == account.id,
            ReportTemplate.project_id.is_(None),
            ReportTemplate.is_active.is_(True),
        ).one_or_none()
        assert template is not None, "Expected a report template for the Trimble project"
        assert template.file_path.startswith("database://")
        assert template.content_bytes
        assert Presentation(io.BytesIO(template.content_bytes)).slides
