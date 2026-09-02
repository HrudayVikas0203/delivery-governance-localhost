from datetime import date, datetime, timedelta, timezone
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import get_settings
from app.core.security import hash_password
from app.db.session import Base, SessionLocal, engine
from app.models.brd import BRDDesignArtifact, BRDDocument, BRDDocumentStatus, BRDRequirementSet
from app.models.delivery import Account, AccountStatus, AllocationRole, Health, Project, ProjectPhase, ResourceAllocation, RiskLevel
from app.models.email import EmailStatus, ScheduledEmail
from app.models.people import Employee, Role
from app.models.status import AIInsight, GeneratedReport, ReportFormat, ReportType, ReportTemplate, SubmissionStatus, WeeklyStatus
from app.models.tasks import Task, TaskAssignment, TaskPriority, TaskStatus
from app.reports.generator import generate_report_file
from app.services.template_storage import store_account_template, validate_pptx_upload


DEMO_PASSWORD = "Demo@123"


def _employee(db, email: str, **values) -> Employee:
    employee = db.query(Employee).filter(Employee.email == email).one_or_none()
    if employee:
        for key, value in values.items():
            setattr(employee, key, value)
        return employee
    employee = Employee(email=email, password_hash=hash_password(DEMO_PASSWORD), **values)
    db.add(employee)
    db.flush()
    return employee


def _account(db, name: str, **values) -> Account:
    account = db.query(Account).filter(Account.name == name).one_or_none()
    if account:
        return account
    account = Account(name=name, **values)
    db.add(account)
    db.flush()
    return account


def _project(db, name: str, account_id: str, **values) -> Project:
    project = db.query(Project).filter(Project.name == name, Project.account_id == account_id).one_or_none()
    if project:
        return project
    project = Project(name=name, account_id=account_id, **values)
    db.add(project)
    db.flush()
    return project


def _allocation(db, project_id: str, employee_id: str, **values) -> None:
    exists = db.query(ResourceAllocation).filter(
        ResourceAllocation.project_id == project_id,
        ResourceAllocation.employee_id == employee_id,
    ).one_or_none()
    if exists:
        for key, value in values.items():
            setattr(exists, key, value)
    else:
        db.add(ResourceAllocation(project_id=project_id, employee_id=employee_id, **values))


def _task(db, project_id: str, title: str, assignee_ids: list[str], **values) -> Task:
    task = db.query(Task).filter(Task.project_id == project_id, Task.title == title).one_or_none()
    if not task:
        task = Task(project_id=project_id, title=title, assignee_id=assignee_ids[0] if assignee_ids else None, **values)
        db.add(task)
        db.flush()
    for employee_id in assignee_ids:
        exists = db.query(TaskAssignment).filter(TaskAssignment.task_id == task.id, TaskAssignment.employee_id == employee_id).one_or_none()
        if not exists:
            db.add(TaskAssignment(task_id=task.id, employee_id=employee_id))
    return task


def _weekly_status(db, employee_id: str, project_id: str, week_start: date, **values) -> None:
    exists = db.query(WeeklyStatus).filter(
        WeeklyStatus.employee_id == employee_id,
        WeeklyStatus.project_id == project_id,
        WeeklyStatus.week_start == week_start,
    ).one_or_none()
    if not exists:
        db.add(WeeklyStatus(employee_id=employee_id, project_id=project_id, week_start=week_start, **values))


def _brd_document(db, project_id: str, filename: str, uploaded_by_id: str, extracted_text: str) -> BRDDocument:
    document = db.query(BRDDocument).filter(BRDDocument.project_id == project_id, BRDDocument.filename == filename).one_or_none()
    if document:
        return document
    document = BRDDocument(
        project_id=project_id,
        filename=filename,
        document_type="brd",
        content_type="text/plain",
        size_bytes=len(extracted_text.encode("utf-8")),
        status=BRDDocumentStatus.READY,
        uploaded_by_id=uploaded_by_id,
        extracted_text=extracted_text,
    )
    db.add(document)
    db.flush()
    return document


def _requirements(db, document: BRDDocument) -> None:
    exists = db.query(BRDRequirementSet).filter(BRDRequirementSet.document_id == document.id, BRDRequirementSet.version == 1).one_or_none()
    if not exists:
        db.add(BRDRequirementSet(
            document_id=document.id,
            project_id=document.project_id,
            version=1,
            overview="Customer portal modernization with secure onboarding, account dashboard, payments, notifications, and audit reporting.",
            functional_json='["Customer login and MFA", "Account summary dashboard", "Payment initiation workflow", "Admin audit reports"]',
            non_functional_json='["RBAC", "99.9% availability", "PII masking", "Exportable compliance reports"]',
            assumptions_json='["Client identity provider is available", "Payment gateway credentials arrive before Sprint 5"]',
            created_by="Demo Seed",
        ))


def _artifact(db, project_id: str, document_id: str, artifact_type: str, title: str, payload_json: str) -> None:
    exists = db.query(BRDDesignArtifact).filter(
        BRDDesignArtifact.project_id == project_id,
        BRDDesignArtifact.artifact_type == artifact_type,
        BRDDesignArtifact.version == 1,
    ).one_or_none()
    if not exists:
        db.add(BRDDesignArtifact(
            project_id=project_id,
            document_id=document_id,
            artifact_type=artifact_type,
            version=1,
            title=title,
            payload_json=payload_json,
            ai_provider="demo",
            model_used="seeded",
        ))


def _create_trimble_template(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = title_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(13, 26, 52)

    bar = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.38))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(26, 120, 220)
    bar.line.fill.background()

    title_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.2), Inches(1.0))
    title_frame = title_box.text_frame
    title = title_frame.paragraphs[0]
    title.text = "Trimble Finance AI Assistant"
    title.alignment = PP_ALIGN.LEFT
    title.runs[0].font.size = Pt(28)
    title.runs[0].font.bold = True
    title.runs[0].font.color.rgb = RGBColor(255, 255, 255)

    subtitle = title_frame.add_paragraph()
    subtitle.text = "Monthly Delivery Status & Governance Review"
    subtitle.alignment = PP_ALIGN.LEFT
    subtitle.runs[0].font.size = Pt(14)
    subtitle.runs[0].font.color.rgb = RGBColor(189, 214, 255)

    badge = title_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.7), Inches(1.5), Inches(2.6), Inches(0.75))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(31, 92, 179)
    badge.line.color.rgb = RGBColor(31, 92, 179)
    badge_tf = badge.text_frame
    badge_tf.text = "Status Deck"
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge_tf.paragraphs[0].font.bold = True
    badge_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    badge_tf.paragraphs[0].font.size = Pt(16)

    summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
    summary_slide.background.fill.solid()
    summary_slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    summary_box = summary_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.5), Inches(12.35), Inches(0.7))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    summary_box.line.color.rgb = RGBColor(191, 219, 254)
    summary_tf = summary_box.text_frame
    summary_tf.text = "Executive Summary"
    summary_tf.paragraphs[0].font.bold = True
    summary_tf.paragraphs[0].font.size = Pt(24)
    summary_tf.paragraphs[0].font.color.rgb = RGBColor(13, 26, 52)

    summary_text = summary_slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.7), Inches(3.3))
    summary_text.text_frame.text = "The Trimble Finance AI Assistant is progressing through delivery governance checkpoints with strong integration momentum across finance workflows, approval routing, and AI trust governance. The project continues to meet major technical milestones while tracking a small number of ERP field mapping dependencies requiring clarification from finance stakeholders."
    for paragraph in summary_text.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(71, 85, 105)

    for i, metrics in enumerate([("Phase", "Development"), ("Sprint", "8"), ("Completion", "62%"), ("Health", "Green")]):
        card = summary_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + (i % 2) * 5.9), Inches(5.0 + (i // 2) * 0.9), Inches(4.4), Inches(0.75))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(219, 234, 254)
        card.line.color.rgb = RGBColor(147, 197, 253)
        tf = card.text_frame
        tf.text = f"{metrics[0]}\n{metrics[1]}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(13, 26, 52)
        tf.paragraphs[1].font.size = Pt(18)
        tf.paragraphs[1].font.bold = True
        tf.paragraphs[1].font.color.rgb = RGBColor(26, 120, 220)

    prs.save(path)


def seed() -> None:
    settings = get_settings()
    Base.metadata.create_all(bind=engine)
    if not settings.seed_demo_data:
        return

    with SessionLocal() as db:
        studio_head = _employee(db, "praveen.baburaya@delta.com", name="Praveen Kumar Baburaya", title="Studio Head", role=Role.STUDIO_HEAD, department="Delivery Leadership")
        pm1 = _employee(db, "gowtham.rallabandi@delta.com", name="Gowtham Rallabandi", title="Program Manager", role=Role.PROGRAM_MANAGER, department="Program Office", manager_id=studio_head.id)
        pm2 = _employee(db, "rambabu.bagati@delta.com", name="Rambabu Bagati", title="Program Manager", role=Role.PROGRAM_MANAGER, department="Program Office", manager_id=studio_head.id)
        proj_m1 = _employee(db, "shanmukha.rewal@delta.com", name="Shanmukha Rewal", title="Project Manager", role=Role.PROJECT_MANAGER, department="Delivery Operations", manager_id=pm1.id)
        proj_m2 = _employee(db, "amrita.kumari@delta.com", name="Amrita Kumari", title="Project Manager", role=Role.PROJECT_MANAGER, department="Delivery Operations", manager_id=pm1.id)
        lead = _employee(db, "ravi.teja@delta.com", name="Ravi Teja Reddy", title="Team Lead", role=Role.TEAM_LEAD, department="Engineering", manager_id=proj_m1.id)
        architect = _employee(db, "suresh.babu@delta.com", name="Suresh Babu", title="Architect", role=Role.TEAM_LEAD, department="Architecture", manager_id=proj_m1.id)
        sr_dev = _employee(db, "deepak.sharma@delta.com", name="Deepak Sharma", title="Senior Developer", role=Role.DEVELOPER, department="Engineering", manager_id=lead.id, skills="React,FastAPI,MySQL")
        dev = _employee(db, "sneha.patil@delta.com", name="Sneha Patil", title="Developer", role=Role.DEVELOPER, department="Engineering", manager_id=lead.id, skills="React,TypeScript")
        qa = _employee(db, "karthik.venkat@delta.com", name="Karthik Venkat", title="QA Engineer", role=Role.DEVELOPER, department="Quality Assurance", manager_id=lead.id, skills="Automation,Playwright")
        devops = _employee(db, "manoj.kumar@delta.com", name="Manoj Kumar", title="DevOps Engineer", role=Role.DEVELOPER, department="Platform", manager_id=lead.id, skills="Docker,Redis,AWS")
        intern = _employee(db, "aditya.verma@delta.com", name="Aditya Verma", title="Engineering Intern", role=Role.INTERN, department="Engineering", manager_id=dev.id)

        acc1 = _account(db, "Acme Corp", industry="Financial Services", country="United States", business_unit="Banking", contract_value=2500000.0, status=AccountStatus.ACTIVE, health=Health.GREEN, delivery_head_id=studio_head.id, program_manager_id=pm1.id)
        acc2 = _account(db, "Global Tech Solutions", industry="Technology", country="United Kingdom", business_unit="Digital", contract_value=1800000.0, status=AccountStatus.ACTIVE, health=Health.GREEN, delivery_head_id=studio_head.id, program_manager_id=pm2.id)
        acc3 = _account(db, "Meridian Health", industry="Healthcare", country="Canada", business_unit="Health Systems", contract_value=3200000.0, status=AccountStatus.ACTIVE, health=Health.AMBER, delivery_head_id=studio_head.id, program_manager_id=pm2.id)
        trimble_account = _account(db, "Trimble", industry="Industrial Technology", country="United States", business_unit="Finance & Field Operations", contract_value=4200000.0, status=AccountStatus.ACTIVE, health=Health.GREEN, delivery_head_id=studio_head.id, program_manager_id=pm1.id)

        proj1 = _project(db, "Retail Banking Portal", acc1.id, phase=ProjectPhase.DEVELOPMENT, health=Health.AMBER, risk=RiskLevel.MEDIUM, client="Acme Corp", budget_used=1.2, budget_total=2.5, program_manager_id=pm1.id, project_manager_id=proj_m1.id, team_lead_id=architect.id, tech_stack="React,FastAPI,MySQL,ChromaDB,Groq", sprint_number=5, description="Modernizing online retail banking portal and customer dashboard.", start_date=date(2026, 6, 1), completion_percent=72)
        proj2 = _project(db, "Fleet Dispatch Engine", acc2.id, phase=ProjectPhase.BETA_TESTING, health=Health.GREEN, risk=RiskLevel.MEDIUM, client="Global Tech Solutions", budget_used=0.9, budget_total=1.8, program_manager_id=pm2.id, project_manager_id=proj_m2.id, team_lead_id=architect.id, tech_stack="Python,FastAPI,Redis,Docker", sprint_number=7, description="Real-time vehicle dispatch engine and route optimization service.", start_date=date(2026, 5, 15), completion_percent=84)
        proj3 = _project(db, "Patient Records Gateway", acc3.id, phase=ProjectPhase.PLANNING, health=Health.AMBER, risk=RiskLevel.HIGH, client="Meridian Health", budget_used=0.4, budget_total=3.2, program_manager_id=pm2.id, project_manager_id=proj_m2.id, team_lead_id=architect.id, tech_stack="Java,Spring Boot,AWS,FHIR", sprint_number=2, description="HIPAA-compliant EHR gateway and patient health record synchronization.", start_date=date(2026, 7, 10), completion_percent=28)
        trimble_pm = _employee(db, "maria.chen@trimble.com", name="Maria Chen", title="Project Manager", role=Role.PROJECT_MANAGER, department="Program Delivery", manager_id=pm1.id)
        trimble_arch = _employee(db, "david.miles@trimble.com", name="David Miles", title="Technical Architect", role=Role.TEAM_LEAD, department="Architecture", manager_id=trimble_pm.id)
        trimble_frontend = _employee(db, "nina.patel@trimble.com", name="Nina Patel", title="Senior Frontend Developer", role=Role.DEVELOPER, department="Engineering", manager_id=trimble_arch.id, skills="React,TypeScript,Tailwind")
        trimble_backend = _employee(db, "omar.hassan@trimble.com", name="Omar Hassan", title="Backend Engineer", role=Role.DEVELOPER, department="Engineering", manager_id=trimble_arch.id, skills="FastAPI,Python,PostgreSQL")
        trimble_data = _employee(db, "jessie.nguyen@trimble.com", name="Jessie Nguyen", title="Data Engineer", role=Role.DEVELOPER, department="Data & AI", manager_id=trimble_arch.id, skills="Azure OpenAI,SQL,Pipelines")
        trimble_qa = _employee(db, "aisha.khan@trimble.com", name="Aisha Khan", title="QA Engineer", role=Role.DEVELOPER, department="Quality Assurance", manager_id=trimble_arch.id, skills="Playwright,Automation")
        trimble_project = _project(db, "Trimble Finance AI Assistant", trimble_account.id, phase=ProjectPhase.DEVELOPMENT, health=Health.GREEN, risk=RiskLevel.MEDIUM, client="Trimble", budget_used=0.85, budget_total=1.8, program_manager_id=pm1.id, project_manager_id=trimble_pm.id, team_lead_id=trimble_arch.id, tech_stack="React,FastAPI,PostgreSQL,Azure OpenAI,Power BI", sprint_number=8, description="AI-assisted finance assistant for approval routing, cash forecasting, and operational finance insights.", start_date=date(2026, 7, 1), completion_percent=62)

        for project, employee, role in [
            (proj1, sr_dev, AllocationRole.DEVELOPER), (proj1, dev, AllocationRole.DEVELOPER), (proj1, qa, AllocationRole.QA),
            (proj2, devops, AllocationRole.DEVOPS), (proj2, sr_dev, AllocationRole.DEVELOPER), (proj3, intern, AllocationRole.INTERN),
            (trimble_project, trimble_pm, AllocationRole.PROJECT_MANAGER),
            (trimble_project, trimble_arch, AllocationRole.TECHNICAL_ARCHITECT),
            (trimble_project, trimble_frontend, AllocationRole.FRONTEND_ENGINEER),
            (trimble_project, trimble_backend, AllocationRole.BACKEND_ENGINEER),
            (trimble_project, trimble_data, AllocationRole.DATA_ENGINEER),
            (trimble_project, trimble_qa, AllocationRole.TESTING_ENGINEER),
        ]:
            _allocation(db, project.id, employee.id, allocation_role=role, allocation_percent=100, start_date=project.start_date or date.today(), reporting_manager_id=project.project_manager_id, created_by_id=studio_head.id)

        today = date.today()
        week_start = today - timedelta(days=today.weekday())

        _task(
            db,
            trimble_project.id,
            "Finalize finance AI approval matrix and workflow logic",
            [trimble_arch.id, trimble_pm.id],
            description="Align finance workflows, exception routing, and approval matrix with client control policies.",
            status=TaskStatus.IN_PROGRESS,
            priority=TaskPriority.CRITICAL,
            due_date=today + timedelta(days=4),
            estimate_hours=22,
            actual_hours=10,
            labels="finance-ai,governance",
            tags=["approval", "workflow"],
            checklist=[{"label": "Workflow map complete", "done": True}, {"label": "Approval matrix validation", "done": False}],
        )
        _task(
            db,
            trimble_project.id,
            "Build invoice approval dashboard and exception queue",
            [trimble_frontend.id],
            description="Deliver the client-facing approval dashboard for invoice exceptions, summaries, and risk indicators.",
            status=TaskStatus.IN_PROGRESS,
            priority=TaskPriority.HIGH,
            due_date=today + timedelta(days=6),
            estimate_hours=26,
            actual_hours=14,
            labels="frontend,ui",
            tags=["dashboard", "finance"],
            checklist=[{"label": "Wireframe complete", "done": True}, {"label": "KPI cards ready", "done": False}],
        )
        _task(
            db,
            trimble_project.id,
            "Integrate ERP ledger and cash forecast connectors",
            [trimble_backend.id],
            description="Connect ERP ledger and forecast APIs with AI assistant workflow and approval endpoints.",
            status=TaskStatus.IN_PROGRESS,
            priority=TaskPriority.HIGH,
            due_date=today + timedelta(days=7),
            estimate_hours=30,
            actual_hours=16,
            labels="backend,erp",
            tags=["api", "integration"],
            checklist=[{"label": "Data contract signed", "done": True}, {"label": "Connector validation", "done": False}],
        )
        _task(
            db,
            trimble_project.id,
            "Validate AI recommendation confidence and guardrails",
            [trimble_data.id],
            description="Check model recommendations against the finance policy, risk guardrails, and control thresholds.",
            status=TaskStatus.REVIEW,
            priority=TaskPriority.HIGH,
            due_date=today + timedelta(days=5),
            estimate_hours=18,
            actual_hours=9,
            labels="ai,quality",
            tags=["guardrails", "review"],
            checklist=[{"label": "Prompt validation", "done": True}, {"label": "Bias and threshold checks", "done": False}],
        )
        _task(
            db,
            trimble_project.id,
            "Complete regression, UAT prep, and release checklist",
            [trimble_qa.id],
            description="Run regression and UAT readiness checks across approval, exception, and finance summary workflows.",
            status=TaskStatus.TODO,
            priority=TaskPriority.MEDIUM,
            due_date=today + timedelta(days=9),
            estimate_hours=16,
            actual_hours=0,
            labels="qa,uat",
            tags=["regression", "release"],
            checklist=[{"label": "Smoke scripts ready", "done": True}, {"label": "Release notes prepared", "done": False}],
        )

        trimble_template_path = get_settings().templates_dir / "trimble_finance_ai_assistant_template.pptx"
        if not trimble_template_path.exists():
            _create_trimble_template(trimble_template_path)

        for offset, (overall_status, completion) in enumerate([("Green", 58), ("Green", 61), ("Amber", 66), ("Green", 72)]):
            week_date = today - timedelta(weeks=3 - offset)
            week_start_date = week_date - timedelta(days=week_date.weekday())
            achievements = [
                "Completed finance workflow mapping and AI prompt design for approval routing.",
                "Integrated ERP ledger and cash forecast data sources for finance review.",
                "Delivered front-end user flows for invoice approvals and exception queues.",
                "Validated AI recommendations against control rules and finance policy guardrails.",
            ]
            status_fields = {
                "project": "Trimble Finance AI Assistant",
                "account": "Trimble",
                "reportingFrequency": "Weekly",
                "overallStatus": overall_status,
                "completionPercent": completion,
                "hoursWorked": 36 + offset * 6,
                "achievements": achievements[offset],
                "blockers": "No critical blockers. Dependency on final ERP field mapping remains under watch." if overall_status == "Green" else "Final ERP field mapping changes require validation before approval automation goes live.",
                "risks": "Finance data quality must stay consistent across invoice, vendor, and cashflow streams.",
                "nextWeekPlan": "Finalize exception workflow tuning and complete regression coverage for invoicing approvals.",
                "supportRequired": "Need confirmation from Trimble finance stakeholders on final approval matrix and SLA thresholds.",
            }
            _weekly_status(db, trimble_backend.id, trimble_project.id, week_start_date, status=SubmissionStatus.SUBMITTED, submitted_at=datetime.now(timezone.utc), fields=status_fields)
            _weekly_status(db, trimble_frontend.id, trimble_project.id, week_start_date, status=SubmissionStatus.SUBMITTED, submitted_at=datetime.now(timezone.utc), fields={**status_fields, "achievements": "Completed the approval queue UX and exception dashboard refinements for finance users.", "hoursWorked": status_fields["hoursWorked"] - 2})
            _weekly_status(db, trimble_qa.id, trimble_project.id, week_start_date, status=SubmissionStatus.APPROVED if offset >= 2 else SubmissionStatus.SUBMITTED, submitted_at=datetime.now(timezone.utc), fields={**status_fields, "achievements": "Regression suite updated for AI recommendation confidence and approval SLA checks.", "blockers": "None; all critical issues are mitigated before release candidate validation.", "hoursWorked": 34 + offset * 5})

        trimble_template = db.query(ReportTemplate).filter(
            ReportTemplate.account_id == trimble_account.id,
            ReportTemplate.project_id.is_(None),
            ReportTemplate.file_type == "pptx",
        ).one_or_none()
        if trimble_template is None or not trimble_template.content_bytes:
            content = trimble_template_path.read_bytes()
            validated = validate_pptx_upload(
                trimble_template_path.name,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                content,
            )
            trimble_template = store_account_template(db, trimble_account.id, validated, trimble_pm.id)

        _weekly_status(db, sr_dev.id, proj1.id, week_start, status=SubmissionStatus.SUBMITTED, submitted_at=datetime.now(timezone.utc), fields={"achievements": "Completed API integration shell and started payment gateway mapping.", "blockers": "Client credentials pending.", "overallStatus": "Amber", "completionPercent": 72, "hoursWorked": 41})
        _weekly_status(db, qa.id, proj1.id, week_start, status=SubmissionStatus.DRAFT, fields={"achievements": "Regression suite updated.", "blockers": "Waiting for stable UAT build.", "overallStatus": "Green", "completionPercent": 68, "hoursWorked": 38})

        _task(db, proj1.id, "Complete payment gateway credential integration", [sr_dev.id, dev.id], description="Wire client-provided payment credentials into sandbox and production config.", status=TaskStatus.BLOCKED, priority=TaskPriority.CRITICAL, due_date=today + timedelta(days=3), estimate_hours=16, actual_hours=5, labels="integration,client-dependency", tags=["payments", "blocked"], checklist=[{"label": "Receive credentials", "done": False}, {"label": "Validate sandbox payment", "done": False}], blocker_reason="Client has not shared credentials.")
        _task(db, proj1.id, "Review BRD requirements coverage", [architect.id], description="Validate extracted BRD requirements against project scope.", status=TaskStatus.REVIEW, priority=TaskPriority.HIGH, due_date=today + timedelta(days=2), estimate_hours=6, labels="brd,architecture", tags=["brd"], checklist=[{"label": "Functional coverage", "done": True}, {"label": "NFR coverage", "done": False}])
        _task(db, proj2.id, "Prepare route optimization demo", [devops.id], description="Package route optimization demo with Redis-backed dispatch simulation.", status=TaskStatus.IN_PROGRESS, priority=TaskPriority.MEDIUM, due_date=today + timedelta(days=5), estimate_hours=12, labels="demo,redis", tags=["demo"])
        _task(db, proj3.id, "Define FHIR entity model", [intern.id, architect.id], description="Draft patient, encounter, provider, and audit entities for review.", status=TaskStatus.TODO, priority=TaskPriority.HIGH, due_date=today + timedelta(days=8), estimate_hours=10, labels="database,healthcare", tags=["fhir"])

        brd_specs = [
            (
                proj1,
                "retail-banking-portal-brd.txt",
                "The retail banking portal supports secure login, MFA, account summaries, payment initiation, notifications, audit reports, RBAC, and client-ready executive reporting. API integration is blocked until client credentials are received.",
                "Retail Banking",
            ),
            (
                proj2,
                "fleet-dispatch-engine-brd.txt",
                "The fleet dispatch engine provides dispatcher dashboards, vehicle telemetry ingestion, route optimization, SLA alerts, depot capacity views, and operational reporting. GPS vendor APIs and route simulation data are key delivery dependencies.",
                "Fleet Dispatch",
            ),
            (
                proj3,
                "patient-records-gateway-brd.txt",
                "The patient records gateway synchronizes FHIR resources across providers, validates consent, logs access, exposes patient search, and reports data quality. Integration with client identity and EHR sandbox environments is required.",
                "Patient Records",
            ),
        ]
        for project, filename, text, label in brd_specs:
            document = _brd_document(db, project.id, filename, architect.id, text)
            _requirements(db, document)
            _artifact(db, project.id, document.id, "business_flow", f"{label} Business Flow", '{"nodes":[{"id":"intake","label":"Intake"},{"id":"validate","label":"Validate"},{"id":"execute","label":"Execute"},{"id":"report","label":"Report"}],"edges":[{"source":"intake","target":"validate"},{"source":"validate","target":"execute"},{"source":"execute","target":"report"}]}')
            _artifact(db, project.id, document.id, "architecture", f"{label} Solution Architecture", '{"layers":[{"name":"Experience","components":["Role Dashboards","Admin Console"]},{"name":"API","components":["FastAPI Gateway","RBAC"]},{"name":"Data","components":["MySQL","ChromaDB"]}],"decisions":["One project ID","Common AI service"]}')

        if not db.query(AIInsight).filter(AIInsight.project_id == proj1.id, AIInsight.week_start == week_start).one_or_none():
            db.add(AIInsight(project_id=proj1.id, week_start=week_start, provider="demo", model="seeded", executive_summary="Delivery is at medium risk because API credential dependency is blocking payment integration.", risk_level="Medium", recommendations={"items": ["Escalate client credential dependency", "Keep QA regression warm", "Review BRD coverage before Sprint 6"]}, health_score=68, sentiment_score=7))

        if not db.query(ScheduledEmail).filter(ScheduledEmail.subject == "Retail Banking Portal weekly dependency update").one_or_none():
            db.add(ScheduledEmail(sender_id=proj_m1.id, recipients=[lead.email, architect.email], subject="Retail Banking Portal weekly dependency update", body="Please review the credential dependency and task blockers before the governance sync.", email_type="project_update", project_id=proj1.id, scheduled_at=datetime.now(timezone.utc) + timedelta(days=1), status=EmailStatus.SCHEDULED))

        db.commit()


if __name__ == "__main__":
    seed()
