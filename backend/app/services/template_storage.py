from collections.abc import Iterator
from contextlib import contextmanager
from dataclasses import dataclass
from hashlib import sha256
from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile

from pptx import Presentation
from sqlalchemy import select
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.models.delivery import Account
from app.models.status import ReportTemplate


PPTX_CONTENT_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/octet-stream",
    "application/zip",
    "",
}


class TemplateStorageError(RuntimeError):
    client_message = "Unable to retrieve the account PPT template."


class TemplateNotConfiguredError(TemplateStorageError):
    client_message = "Account PPT template is not configured."


class InvalidTemplateError(TemplateStorageError):
    client_message = "The configured account PPT template is invalid or corrupted."


class TemplateUploadError(ValueError):
    pass


@dataclass(frozen=True)
class ValidatedTemplate:
    filename: str
    content_type: str
    content: bytes
    content_sha256: str
    slide_count: int


def _safe_filename(filename: str | None) -> str:
    return Path((filename or "").replace("\\", "/")).name


def validate_pptx_upload(filename: str | None, content_type: str | None, content: bytes) -> ValidatedTemplate:
    safe_filename = _safe_filename(filename)
    if Path(safe_filename).suffix.lower() != ".pptx":
        raise TemplateUploadError("Only .pptx account templates are supported.")
    if (content_type or "").lower() not in PPTX_CONTENT_TYPES:
        raise TemplateUploadError("The uploaded file type is not a supported PPTX MIME type.")
    if not content:
        raise TemplateUploadError("The uploaded PPTX template is empty.")
    if len(content) > get_settings().ppt_template_max_bytes:
        max_mb = get_settings().ppt_template_max_bytes // (1024 * 1024)
        raise TemplateUploadError(f"The PPTX template exceeds the {max_mb} MB size limit.")
    try:
        presentation = Presentation(BytesIO(content))
    except Exception as exc:
        raise TemplateUploadError("The uploaded PPTX template is invalid or corrupted.") from exc
    if not presentation.slides:
        raise TemplateUploadError("The uploaded PPTX template must contain at least one slide.")
    return ValidatedTemplate(
        filename=safe_filename,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        content=content,
        content_sha256=sha256(content).hexdigest(),
        slide_count=len(presentation.slides),
    )


def find_account_template(db: Session, account_id: str, *, lock: bool = False) -> ReportTemplate | None:
    statement = (
        select(ReportTemplate)
        .where(
            ReportTemplate.account_id == account_id,
            ReportTemplate.project_id.is_(None),
            ReportTemplate.file_type == "pptx",
            ReportTemplate.is_active.is_(True),
        )
        .order_by(ReportTemplate.updated_at.desc(), ReportTemplate.uploaded_at.desc())
    )
    if lock:
        statement = statement.with_for_update()
    return db.scalar(statement)


def store_account_template(
    db: Session,
    account_id: str,
    validated: ValidatedTemplate,
    uploaded_by_id: str,
) -> ReportTemplate:
    # Locking the account serializes concurrent replacements for the same account.
    account = db.scalar(select(Account).where(Account.id == account_id).with_for_update())
    if account is None:
        raise LookupError("Account not found")

    template = find_account_template(db, account_id, lock=True)
    if template is None:
        template = ReportTemplate(
            name=Path(validated.filename).stem,
            file_path="database://pending",
            file_type="pptx",
            account_id=account_id,
            project_id=None,
        )
        db.add(template)
        db.flush()

    template.name = Path(validated.filename).stem
    template.filename = validated.filename
    template.file_path = f"database://report-templates/{template.id}"
    template.file_type = "pptx"
    template.content_type = validated.content_type
    template.size_bytes = len(validated.content)
    template.content_bytes = validated.content
    template.content_sha256 = validated.content_sha256
    template.uploaded_by_id = uploaded_by_id
    template.project_id = None
    template.account_id = account_id
    template.is_active = True

    # Preserve legacy duplicate metadata, but prevent it from being selected.
    duplicates = db.scalars(
        select(ReportTemplate).where(
            ReportTemplate.account_id == account_id,
            ReportTemplate.project_id.is_(None),
            ReportTemplate.file_type == "pptx",
            ReportTemplate.id != template.id,
            ReportTemplate.is_active.is_(True),
        )
    ).all()
    for duplicate in duplicates:
        duplicate.is_active = False
    return template


def remove_account_templates(db: Session, account_id: str) -> list[ReportTemplate]:
    templates = list(
        db.scalars(
            select(ReportTemplate).where(
                ReportTemplate.account_id == account_id,
                ReportTemplate.project_id.is_(None),
                ReportTemplate.file_type == "pptx",
            )
        ).all()
    )
    for template in templates:
        db.delete(template)
    return templates


def get_account_template_bytes(db: Session, account_id: str, template_id: str) -> bytes:
    template = db.get(ReportTemplate, template_id)
    if (
        template is None
        or template.account_id != account_id
        or template.project_id is not None
        or template.file_type != "pptx"
        or not template.is_active
    ):
        raise TemplateNotConfiguredError()
    if not template.content_bytes:
        raise TemplateStorageError()
    content = bytes(template.content_bytes)
    if template.content_sha256 and sha256(content).hexdigest() != template.content_sha256:
        raise InvalidTemplateError()
    return content


@contextmanager
def get_account_template_file(db: Session, account_id: str, template_id: str) -> Iterator[Path]:
    """Materialize database-authoritative template bytes into a unique temporary file."""
    content = get_account_template_bytes(db, account_id, template_id)
    temporary_path: Path | None = None
    try:
        with NamedTemporaryFile(prefix="account-ppt-template-", suffix=".pptx", delete=False) as temporary:
            temporary.write(content)
            temporary.flush()
            temporary_path = Path(temporary.name)
        try:
            Presentation(str(temporary_path))
        except Exception as exc:
            raise InvalidTemplateError() from exc
        yield temporary_path
    finally:
        if temporary_path is not None:
            temporary_path.unlink(missing_ok=True)
