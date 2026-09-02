import json
import logging
import re
from pathlib import Path
from textwrap import shorten
import time

from openpyxl import Workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from sqlalchemy.orm import Session
from pydantic import BaseModel, Field

from app.core.config import get_settings
from app.ai.ppt_mapping import GeminiMappingConfigurationError, GeminiMappingError, map_template
from app.models.delivery import Account, Project
from app.models.people import Employee
from app.models.status import GeneratedReport, ReportFormat, WeeklyStatus
from app.models.tasks import Task, TaskStatus
from app.schemas.common import LLMSelection
from app.services.template_storage import TemplateNotConfiguredError, get_account_template_file


BRAND_BLUE = "#1D4ED8"
BRAND_NAVY = "#0F172A"
INK_SOFT = "#475569"
SUCCESS = "#16A34A"
WARNING = "#D97706"
DANGER = "#DC2626"
LIGHT_BLUE = "#EFF6FF"
BORDER_BLUE = "#BFDBFE"
logger = logging.getLogger(__name__)


class ReportNarrative(BaseModel):
    executive_summary: str = Field(min_length=1)
    achievements: list[str] = Field(default_factory=list)
    risks: list[str] = Field(default_factory=list)
    blockers: list[str] = Field(default_factory=list)
    next_steps: list[str] = Field(default_factory=list)


def _report_path(report: GeneratedReport) -> Path:
    suffix = {"pptx": ".pptx", "pdf": ".pdf", "xlsx": ".xlsx"}[report.report_format.value]
    safe_title = "".join(ch if ch.isalnum() else "_" for ch in report.title.lower()).strip("_")
    return get_settings().report_dir / f"{report.id}_{safe_title}{suffix}"


def _project_rows(db: Session, report: GeneratedReport) -> list[Project]:
    query = db.query(Project)
    if "project:" in report.scope:
        project_id = report.scope.split("project:", 1)[1].strip().split()[0]
        project = db.get(Project, project_id)
        return [project] if project else []
    if "account:" in report.scope:
        account_id = report.scope.split("account:", 1)[1].strip().split()[0]
        query = query.filter(Project.account_id == account_id)
    return query.order_by(Project.name).all()


def _scope_value(scope: str, key: str) -> str | None:
    marker = f"{key}:"
    if marker not in scope:
        return None
    return scope.split(marker, 1)[1].strip().split()[0]


def _status_rows(db: Session, report: GeneratedReport, projects: list[Project]) -> list[WeeklyStatus]:
    query = db.query(WeeklyStatus)
    employee_id = _scope_value(report.scope, "employee")
    project_id = _scope_value(report.scope, "project")
    account_id = _scope_value(report.scope, "account")
    period = _scope_value(report.scope, "period")

    if employee_id:
        query = query.filter(WeeklyStatus.employee_id == employee_id)
    if project_id:
        query = query.filter(WeeklyStatus.project_id == project_id)
    elif account_id:
        project_ids = [project.id for project in projects]
        if project_ids:
            query = query.filter(WeeklyStatus.project_id.in_(project_ids))
    rows = query.order_by(WeeklyStatus.week_start.desc()).limit(120).all()
    if period:
        labels = {"daily": "Daily", "weekly": "Weekly", "monthly": "Monthly"}
        expected = labels.get(period.lower())
        if expected:
            rows = [
                row for row in rows
                if str(row.fields.get("reportingFrequency") or row.fields.get("frequency") or "Weekly") == expected
            ]
    return rows[:60]


def _text(value: object | None, fallback: str = "Not reported") -> str:
    if value is None:
        return fallback
    cleaned = str(value).strip()
    return cleaned or fallback


def _short(value: object | None, limit: int = 120, fallback: str = "Not reported") -> str:
    return shorten(_text(value, fallback), width=limit, placeholder="...")


def _clean_llm_text(value: str) -> str:
    cleaned = value.replace("**", "").replace("__", "").replace("###", "").strip()
    return "\n".join(line.strip(" -") for line in cleaned.splitlines() if line.strip())


def _is_meaningful(value: object | None) -> bool:
    if value is None:
        return False
    cleaned = str(value).strip().lower()
    return cleaned not in {"", "none", "no", "no blockers", "n/a", "na", "not applicable"}


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _set_slide_title(slide, title: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = title
        title_shape = slide.shapes.title
    else:
        title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(9.1), Inches(0.55))
        title_shape.text_frame.text = title
    p = title_shape.text_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = _rgb(BRAND_NAVY)


def _add_textbox(slide, text: str, x: float, y: float, w: float, h: float, size: int = 12, bold: bool = False, color: str = INK_SOFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.text = text
    for paragraph in frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = _rgb(color)


def _add_metric_card(slide, label: str, value: str | int, x: float, y: float, w: float = 1.65, h: float = 0.9, accent: str = BRAND_BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(LIGHT_BLUE)
    shape.line.color.rgb = _rgb(BORDER_BLUE)
    tf = shape.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.text = str(value)
    tf.paragraphs[0].font.size = Pt(19)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _rgb(accent)
    p = tf.add_paragraph()
    p.text = label
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = _rgb(INK_SOFT)


def _metrics(statuses: list[WeeklyStatus], projects: list[Project]) -> dict:
    submitted = [s for s in statuses if s.status.value in {"submitted", "approved"}]
    approved = [s for s in statuses if s.status.value == "approved"]
    blockers = [s for s in statuses if _is_meaningful(s.fields.get("blockers")) or str(s.fields.get("overallStatus", "")).lower() == "red"]
    hours = sum(int(s.fields.get("hoursWorked") or 0) for s in statuses)
    completion_values = [int(s.fields.get("completionPercent") or 0) for s in statuses if s.fields.get("completionPercent") is not None]
    avg_completion = round(sum(completion_values) / len(completion_values)) if completion_values else round(sum(p.completion_percent for p in projects) / len(projects)) if projects else 0
    health_counts = {"Green": 0, "Amber": 0, "Red": 0}
    for status in statuses:
        health = str(status.fields.get("overallStatus") or "Green")
        if health in health_counts:
            health_counts[health] += 1
    return {
        "project_count": len(projects),
        "status_count": len(statuses),
        "submitted_count": len(submitted),
        "approved_count": len(approved),
        "blocker_count": len(blockers),
        "hours": hours,
        "avg_completion": avg_completion,
        "health_counts": health_counts,
        "blockers": blockers,
    }


def _fallback_summary(report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus]) -> str:
    metrics = _metrics(statuses, projects)
    high_risk_projects = [project.name for project in projects if project.risk.value in {"high", "critical"} or project.health.value == "red"]
    risk_phrase = "; ".join(high_risk_projects[:3]) if high_risk_projects else "no critical project-level risk currently flagged"
    return (
        f"{report.title} covers {metrics['project_count']} project(s) and {metrics['status_count']} status update(s). "
        f"Average completion is {metrics['avg_completion']}%, with {metrics['approved_count']} approved status update(s) and "
        f"{metrics['blocker_count']} blocker or red-health item(s). Current governance attention should focus on {risk_phrase}. "
        "The recommended leadership action is to close blockers, confirm dependencies, and keep weekly approvals current."
    )


def _llm_summary(report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus], llm: LLMSelection | None) -> str:
    fallback = _fallback_summary(report, projects, statuses)
    if not llm:
        from app.schemas.common import LLMSelection
        from app.core.config import get_settings
        llm = LLMSelection(provider="gemini", model=get_settings().gemini_default_model)
    from app.services.llm import generate_text

    status_digest = "\n".join(
        f"- {_text(status.fields.get('project'), 'Project')} | {status.week_start} | {status.status.value} | "
        f"health={_text(status.fields.get('overallStatus'), 'Green')} | complete={status.fields.get('completionPercent') or 0}% | "
        f"period={_text(status.fields.get('reportingFrequency') or status.fields.get('frequency'), 'Weekly')} | "
        f"update={_short(status.fields.get('achievements'), 140)} | blocker={_short(status.fields.get('blockers'), 100, 'None')}"
        for status in statuses[:16]
    )
    metrics = _metrics(statuses, projects)
    prompt = (
        "Return only valid JSON matching this schema: {\"executive_summary\": string, \"achievements\": [string], \"risks\": [string], \"blockers\": [string], \"next_steps\": [string]}. "
        "Write premium, client-ready PowerPoint copy in 120-150 words for a delivery governance status deck. "
        "Use the exact evidence below; do not invent progress, dates, people, budget, or milestones. "
        "Structure the copy as three concise labeled lines: Delivery position, Risk posture, Leadership actions. "
        "Make the language executive, specific, and analytical. Avoid filler, apologies, markdown, asterisks, and generic phrases such as risks and challenges. "
        "Do not describe contributor status updates as separate projects; if there is one project, say one project with multiple contributors.\n\n"
        f"Report title: {report.title}\nScope: {report.scope}\nRequested period filter: {_scope_value(report.scope, 'period') or 'all'}\n"
        f"Project count: {metrics['project_count']}\nStatus update count: {metrics['status_count']}\n"
        f"Average completion: {metrics['avg_completion']}%\nBlocker/red item count: {metrics['blocker_count']}\n"
        f"Health distribution: {metrics['health_counts']}\n"
        f"Projects: {', '.join(project.name for project in projects) if projects else 'None'}\n"
        f"Contributor status digest:\n{status_digest or 'No status updates available.'}"
    )
    try:
        text, _ = generate_text(llm.provider, prompt, llm.model)
        narrative = ReportNarrative.model_validate(json.loads(text))
        return _clean_llm_text(narrative.executive_summary)
    except Exception:
        return fallback


def _template_values(report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus], db: Session, summary: str) -> dict[str, object]:
    account_id = _scope_value(report.scope, "account")
    account = db.get(Account, account_id) if account_id else None
    metrics = _metrics(statuses, projects)
    achievements = list(dict.fromkeys(
        _short(status.fields.get("achievements"), 180)
        for status in statuses
        if _is_meaningful(status.fields.get("achievements"))
    ))
    risks = list(dict.fromkeys(
        _short(status.fields.get("risks"), 180)
        for status in statuses
        if _is_meaningful(status.fields.get("risks"))
    ))
    blockers = list(dict.fromkeys(
        _short(status.fields.get("blockers"), 180)
        for status in statuses
        if _is_meaningful(status.fields.get("blockers"))
    ))
    next_steps = list(dict.fromkeys(
        _short(status.fields.get("nextWeekPlan"), 180)
        for status in statuses
        if _is_meaningful(status.fields.get("nextWeekPlan"))
    ))
    overall_status = "Red" if metrics["health_counts"]["Red"] else "Amber" if metrics["health_counts"]["Amber"] else "Green"
    project_ids = [project.id for project in projects]
    tasks = db.query(Task).filter(Task.project_id.in_(project_ids)).order_by(Task.updated_at.desc()).all() if project_ids else []
    completed_tasks = sum(task.status == TaskStatus.DONE for task in tasks)
    blocked_tasks = sum(task.status == TaskStatus.BLOCKED for task in tasks)
    coverage = None
    coverage_metrics: dict = {}
    try:
        from app.services.code_quality import get_coverage
        coverage_payload = get_coverage()
        coverage = coverage_payload.get("overall")
        coverage_metrics = coverage_payload.get("metrics") or {}
    except Exception:
        logger.warning("Coverage data was unavailable for report %s", report.id)
    bullet = lambda items, fallback: "\n".join(f"• {item}" for item in items[:5]) if items else fallback
    task_rows = [[task.title, task.status.value.replace("_", " ").title(), task.priority.value.title()] for task in tasks[:6]]
    return {
        "TITLE": report.title,
        "REPORT_TITLE": report.title,
        "SUMMARY": summary,
        "PROJECT": ", ".join(project.name for project in projects) or "All Projects",
        "PROJECT_NAME": ", ".join(project.name for project in projects) or "All Projects",
        "DATE": report.generated_at.strftime("%d %b %Y") if report.generated_at else "",
        "ACCOUNT_NAME": account.name if account else "All Accounts",
        "PROJECT_NAME": ", ".join(project.name for project in projects) or "All Projects",
        "REPORT_DATE": report.generated_at.strftime("%d %b %Y") if report.generated_at else "",
        "OVERALL_STATUS": overall_status,
        "COMPLETION": f"{metrics['avg_completion']}%",
        "HOURS": str(metrics["hours"]),
        "EXECUTIVE_SUMMARY": summary,
        "ACHIEVEMENTS": "\n".join(f"- {item}" for item in achievements[:8]) or "No achievements reported",
        "RISKS": "\n".join(f"- {item}" for item in risks[:8]) or "No risks reported",
        "BLOCKERS": "\n".join(f"- {item}" for item in blockers[:8]) or "No blockers reported",
        "NEXT_STEPS": "\n".join(f"- {item}" for item in next_steps[:8]) or "No next steps reported",
        "NEXT_WEEK_PLAN": "\n".join(f"- {item}" for item in next_steps[:8]) or "No next steps reported",
        "PROJECT_METRICS": f"Projects: {metrics['project_count']} | Updates: {metrics['status_count']} | Completion: {metrics['avg_completion']}% | Blockers: {metrics['blocker_count']}",
        "CODE_COVERAGE": f"{coverage:.1f}%" if isinstance(coverage, (int, float)) else "No data available",
        "TASKS": f"{len(tasks)} Tasks · {completed_tasks} Completed · {blocked_tasks} Blocked" if tasks else "No data available",
        "TASKS_COMPLETED": f"{completed_tasks} Completed" if tasks else "No data available",
        "KEY_RISKS": bullet(risks, "No data available"),
        "TASKS_TABLE": task_rows,
        "TASKS_CHART": {"categories": ["Completed", "In Progress", "Review", "Blocked", "To Do"], "values": [completed_tasks, sum(task.status == TaskStatus.IN_PROGRESS for task in tasks), sum(task.status == TaskStatus.REVIEW for task in tasks), blocked_tasks, sum(task.status == TaskStatus.TODO for task in tasks)]},
        "COVERAGE_CHART": {"categories": ["Statements", "Branches", "Functions", "Lines"], "values": [coverage_metrics.get(key, 0) or 0 for key in ("statements", "branches", "functions", "lines")]},
    }


def _populate_template(prs: Presentation, values: dict[str, object], mapping=None) -> bool:
    replaced = False
    normalized_values = {key.lower(): str(value) for key, value in values.items() if isinstance(value, (str, int, float))}

    def placeholder_key(text: str) -> str | None:
        match = re.fullmatch(r"\s*\{\{\s*([A-Za-z0-9_]+)\s*\}\}\s*", text)
        return match.group(1).upper() if match else None

    def add_table(slide, shape, rows: list[list[str]]) -> None:
        headers = ["Task", "Status", "Priority"]
        table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), shape.left, shape.top, shape.width, shape.height)
        table = table_shape.table
        for col_index, header in enumerate(headers):
            cell = table.cell(0, col_index)
            cell.text = header
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(BRAND_NAVY)
            cell.text_frame.paragraphs[0].font.bold = True; cell.text_frame.paragraphs[0].font.color.rgb = _rgb("FFFFFF")
        for row_index, row in enumerate(rows, start=1):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = _short(value, 65, "No data available")
                cell.text_frame.word_wrap = True
                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].font.size = Pt(9)
        table.columns[0].width = int(shape.width * 0.58)
        table.columns[1].width = int(shape.width * 0.25)
        table.columns[2].width = int(shape.width * 0.17)

    def add_chart(slide, shape, chart_values: dict) -> None:
        categories = chart_values.get("categories") or []
        values = chart_values.get("values") or []
        if not categories or not values:
            return
        max_value = max(max(values), 1)
        padding = int(shape.width * 0.02)
        label_width = int(shape.width * 0.26)
        value_width = int(shape.width * 0.10)
        bar_width = shape.width - label_width - value_width - padding * 3
        row_height = max(int(shape.height / len(categories)), 1)
        for index, (category, value) in enumerate(zip(categories, values)):
            top = shape.top + row_height * index + int(row_height * 0.18)
            height = int(row_height * 0.55)
            label = slide.shapes.add_textbox(shape.left + padding, top, label_width - padding, height)
            label.text_frame.text = str(category)
            label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            label.text_frame.paragraphs[0].font.size = Pt(9)
            track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, shape.left + label_width + padding, top, bar_width, height)
            track.fill.solid(); track.fill.fore_color.rgb = _rgb("E2E8F0"); track.line.fill.background()
            actual = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, shape.left + label_width + padding, top, max(1, int(bar_width * float(value) / max_value)), height)
            actual.fill.solid(); actual.fill.fore_color.rgb = _rgb(BRAND_BLUE); actual.line.fill.background()
            value_label = slide.shapes.add_textbox(shape.left + label_width + bar_width + padding * 2, top, value_width, height)
            value_label.text_frame.text = f"{float(value):.1f}%" if categories[0] == "Statements" else str(int(value))
            value_label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            value_label.text_frame.paragraphs[0].font.size = Pt(9)

    def replace_text(text: str) -> str:
        return re.sub(
            r"\{\{\s*([A-Za-z0-9_]+)\s*\}\}|\[\s*([A-Za-z0-9_]+)\s*\]",
            lambda match: normalized_values.get((match.group(1) or match.group(2)).lower(), match.group(0)),
            text,
        )

    def set_text_frame(text_frame, text: str) -> None:
        nonlocal replaced
        if text == text_frame.text:
            return
        paragraphs = text_frame.paragraphs
        first_paragraph = paragraphs[0]
        if first_paragraph.runs:
            first_paragraph.runs[0].text = text
            for run in first_paragraph.runs[1:]:
                run.text = ""
        else:
            first_paragraph.add_run().text = text
        for paragraph in paragraphs[1:]:
            for run in paragraph.runs:
                run.text = ""
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        replaced = True

    def replace_frame_tokens(text_frame) -> None:
        nonlocal replaced
        # PowerPoint may split a token across runs. Update the inherited first
        # run so the template's existing shape, dimensions, and style remain.
        full_text = text_frame.text
        updated_full_text = replace_text(full_text)
        if updated_full_text != full_text:
            set_text_frame(text_frame, updated_full_text)
            return
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                updated = replace_text(run.text)
                if updated != run.text:
                    run.text = updated
                    replaced = True

    def mapped_text(fields: list[str]) -> str:
        parts = []
        for key in fields:
            source_value = values.get(key)
            if source_value is not None:
                parts.append(str(source_value))
        return "\n".join(parts)

    mapped_elements = {}
    for slide_mapping in (mapping.slides if mapping else []):
        for element_id, content in slide_mapping.mapped_content.items():
            mapped_elements[element_id] = content

    def populate_shape(slide, shape, element_id: str) -> None:
            nonlocal replaced
            if getattr(shape, "has_text_frame", False):
                key = placeholder_key(shape.text_frame.text)
                rich_value = values.get(key) if key else None
                if isinstance(rich_value, list):
                    add_table(slide, shape, rich_value)
                    shape._element.getparent().remove(shape._element)
                    replaced = True
                    return
                if isinstance(rich_value, dict):
                    add_chart(slide, shape, rich_value)
                    shape._element.getparent().remove(shape._element)
                    replaced = True
                    return
                if element_id in mapped_elements:
                    set_text_frame(shape.text_frame, mapped_elements[element_id])
                else:
                    replace_frame_tokens(shape.text_frame)
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for cell_index, cell in enumerate(row.cells):
                        cell_id = f"{element_id}_cell_{row_index}_{cell_index}"
                        if cell_id in mapped_elements:
                            set_text_frame(cell.text_frame, mapped_elements[cell_id])
                        else:
                            replace_frame_tokens(cell.text_frame)
            if hasattr(shape, "shapes"):
                for child_index, child in enumerate(shape.shapes):
                    populate_shape(slide, child, f"{element_id}_group_{child_index}")

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            populate_shape(slide, shape, f"slide_{slide_index}_shape_{shape_index}")
    return replaced


def generate_report_file(db: Session, report_id: str, llm: LLMSelection | None = None) -> str:
    report = db.get(GeneratedReport, report_id)
    if report is None:
        raise ValueError("Report not found")
    projects = _project_rows(db, report)
    statuses = _status_rows(db, report, projects)
    path = _report_path(report)
    account_id = _scope_value(report.scope, "account")
    project_ids = [project.id for project in projects]
    report_started = time.perf_counter()
    stage = "report_generation"
    logger.info("REPORT_GENERATION_BEGIN report_id=%s account_id=%s project_ids=%s stage=%s elapsed_ms=0", report.id, account_id, project_ids, stage)
    try:
        if report.report_format == ReportFormat.PPTX:
            account_template = (
                report.template
                if account_id
                and getattr(report, "template", None)
                and report.template.account_id == account_id
                and report.template.project_id is None
                and report.template.file_type == "pptx"
                and report.template.is_active
                else None
            )
            if not account_id or account_template is None:
                raise TemplateNotConfiguredError()
            stage = "template_lookup"
            stage_started = time.perf_counter()
            logger.info("TEMPLATE_LOOKUP_BEGIN report_id=%s account_id=%s project_ids=%s stage=%s elapsed_ms=0", report.id, account_id, project_ids, stage)
            with get_account_template_file(db, account_id, account_template.id) as template_path:
                logger.info("TEMPLATE_LOOKUP_END report_id=%s account_id=%s project_ids=%s stage=%s elapsed_ms=%d", report.id, account_id, project_ids, stage, round((time.perf_counter() - stage_started) * 1000))
                _generate_ppt(path, report, projects, statuses, db, llm, template_path, account_id)
        elif report.report_format == ReportFormat.PDF:
            _generate_pdf(path, report, projects, statuses, db, llm)
        else:
            _generate_xlsx(path, report, projects, statuses, db)

        content = path.read_bytes()
        content_types = {
            ReportFormat.PPTX: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ReportFormat.PDF: "application/pdf",
            ReportFormat.XLSX: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        }
        report.file_path = str(path)
        report.filename = path.name
        report.content_type = content_types[report.report_format]
        report.size_bytes = len(content)
        report.content_bytes = content
        report.status = "ready"
        db.commit()
        logger.info("REPORT_GENERATION_END report_id=%s account_id=%s project_ids=%s stage=report_generation elapsed_ms=%d", report.id, account_id, project_ids, round((time.perf_counter() - report_started) * 1000))
        return str(path)
    except Exception as exc:
        logger.exception(
            "REPORT_GENERATION_FAILED report_id=%s account_id=%s project_ids=%s stage=%s exception_type=%s elapsed_ms=%d",
            report.id,
            account_id,
            project_ids,
            stage,
            type(exc).__name__,
            round((time.perf_counter() - report_started) * 1000),
        )
        raise


def _generate_ppt(
    path: Path,
    report: GeneratedReport,
    projects: list[Project],
    statuses: list[WeeklyStatus],
    db: Session,
    llm: LLMSelection | None = None,
    template_path: Path | None = None,
    account_id: str | None = None,
) -> None:
    if template_path is None:
        raise TemplateNotConfiguredError()

    project_ids = [project.id for project in projects]
    validation_started = time.perf_counter()
    logger.info("TEMPLATE_VALIDATION_BEGIN report_id=%s account_id=%s project_ids=%s stage=template_validation elapsed_ms=0", report.id, account_id, project_ids)
    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        from app.services.template_storage import InvalidTemplateError

        raise InvalidTemplateError() from exc
    if not prs.slides:
        from app.services.template_storage import InvalidTemplateError

        raise InvalidTemplateError()
    original_slide_count = len(prs.slides)
    logger.info("TEMPLATE_VALIDATION_END report_id=%s account_id=%s project_ids=%s stage=template_validation elapsed_ms=%d", report.id, account_id, project_ids, round((time.perf_counter() - validation_started) * 1000))

    summary_text = _llm_summary(report, projects, statuses, llm)
    values = _template_values(report, projects, statuses, db, summary_text)

    # Tokenized templates are fully functional without an AI integration.
    # Gemini mapping remains optional for explicitly configured element maps.
    mapping_started = time.perf_counter()
    logger.info("GEMINI_MAPPING_BEGIN report_id=%s account_id=%s project_ids=%s stage=gemini_mapping elapsed_ms=0", report.id, account_id, project_ids)
    _, mapping = map_template(
        template_path,
        db,
        projects,
        statuses,
        {"id": report.id, "title": report.title, "type": report.report_type.value, "scope": report.scope},
    )
    logger.info("GEMINI_MAPPING_END report_id=%s account_id=%s project_ids=%s stage=gemini_mapping elapsed_ms=%d", report.id, account_id, project_ids, round((time.perf_counter() - mapping_started) * 1000))

    population_started = time.perf_counter()
    logger.info("PPT_POPULATION_BEGIN report_id=%s account_id=%s project_ids=%s stage=ppt_population elapsed_ms=0", report.id, account_id, project_ids)
    _populate_template(prs, values, mapping)
    prs.save(path)
    try:
        generated = Presentation(str(path))
    except Exception as exc:
        raise RuntimeError("Generated PPT could not be reopened.") from exc
    if len(generated.slides) != original_slide_count:
        raise RuntimeError("Generated PPT changed the uploaded template slide structure.")
    logger.info("PPT_POPULATION_END report_id=%s account_id=%s project_ids=%s stage=ppt_population elapsed_ms=%d", report.id, account_id, project_ids, round((time.perf_counter() - population_started) * 1000))


def _generate_pdf(path: Path, report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus], db: Session, llm: LLMSelection | None = None) -> None:
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(path), pagesize=A4, rightMargin=28, leftMargin=28, topMargin=28, bottomMargin=28)
    metrics = _metrics(statuses, projects)
    summary_text = _llm_summary(report, projects, statuses, llm)
    story = [Paragraph(report.title, styles["Title"]), Paragraph("AI-assisted delivery governance report", styles["Normal"]), Spacer(1, 12), Paragraph(summary_text, styles["BodyText"]), Spacer(1, 14)]

    kpi_data = [["Projects", "Updates", "Approved", "Blocked", "Avg Complete"], [metrics["project_count"], metrics["status_count"], metrics["approved_count"], metrics["blocker_count"], f"{metrics['avg_completion']}%"]]
    kpi_table = Table(kpi_data, repeatRows=1)
    kpi_table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1D4ED8")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("ALIGN", (0, 0), (-1, -1), "CENTER"), ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#BFDBFE")), ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold")]))
    story.extend([kpi_table, Spacer(1, 16)])

    data = [["Project", "Account", "Phase", "Health", "Risk", "Completion"]]
    for project in projects:
        account = db.get(Account, project.account_id)
        data.append([project.name, account.name if account else "-", project.phase.value, project.health.value, project.risk.value, f"{project.completion_percent}%"])
    table = Table(data, repeatRows=1)
    table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0F172A")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("GRID", (0, 0), (-1, -1), 0.25, colors.grey), ("VALIGN", (0, 0), (-1, -1), "TOP")]))
    story.append(table)

    if statuses:
        story.extend([Spacer(1, 18), Paragraph("Status Details", styles["Heading2"])])
        status_data = [["Person", "Cycle", "Status", "Health", "Key Update"]]
        for status in statuses[:14]:
            employee = db.get(Employee, status.employee_id)
            status_data.append([employee.name if employee else "-", f"{_text(status.fields.get('reportingFrequency') or status.fields.get('frequency'), 'Weekly')} / {status.week_start.strftime('%d %b %Y')}", status.status.value, _text(status.fields.get("overallStatus"), "-"), _short(status.fields.get("achievements"), 95, "No update")])
        status_table = Table(status_data, repeatRows=1)
        status_table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563EB")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("GRID", (0, 0), (-1, -1), 0.25, colors.grey), ("VALIGN", (0, 0), (-1, -1), "TOP")]))
        story.append(status_table)
    doc.build(story)

    if getattr(report, "template", None) and report.template and report.template.file_type == "pdf":
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError as exc:
            raise RuntimeError("PDF template support requires pypdf to be installed") from exc
        final_writer = PdfWriter()
        template_reader = PdfReader(report.template.file_path)
        output_reader = PdfReader(str(path))
        for page in template_reader.pages:
            final_writer.add_page(page)
        for page in output_reader.pages:
            final_writer.add_page(page)
        with path.open("wb") as output_file:
            final_writer.write(output_file)


def _generate_xlsx(path: Path, report: GeneratedReport, projects: list[Project], statuses: list[WeeklyStatus], db: Session) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Delivery Governance"
    ws.append(["Project", "Account", "Phase", "Health", "Risk", "Budget Used", "Budget Total", "Completion"])
    for project in projects:
        account = db.get(Account, project.account_id)
        ws.append([project.name, account.name if account else "-", project.phase.value, project.health.value, project.risk.value, float(project.budget_used), float(project.budget_total), project.completion_percent])
    for col in ws.columns:
        ws.column_dimensions[col[0].column_letter].width = 18
    status_ws = wb.create_sheet("Status Details")
    status_ws.append(["Employee", "Project", "Reporting Cycle", "Period Start", "Submission Status", "Health", "Completion", "Hours", "Achievements", "Blockers", "Risks", "Next Plan"])
    for status in statuses:
        employee = db.get(Employee, status.employee_id)
        project = db.get(Project, status.project_id) if status.project_id else None
        status_ws.append([employee.name if employee else "-", project.name if project else _text(status.fields.get("project"), "-"), _text(status.fields.get("reportingFrequency") or status.fields.get("frequency"), "Weekly"), status.week_start.isoformat(), status.status.value, _text(status.fields.get("overallStatus"), "-"), status.fields.get("completionPercent") or 0, status.fields.get("hoursWorked") or 0, _text(status.fields.get("achievements"), ""), _text(status.fields.get("blockers"), ""), _text(status.fields.get("risks"), ""), _text(status.fields.get("nextWeekPlan"), "")])
    for col in status_ws.columns:
        status_ws.column_dimensions[col[0].column_letter].width = 22
    wb.save(path)
