from pathlib import Path

from pydantic import BaseModel, Field
from pptx import Presentation


class TemplateElement(BaseModel):
    id: str
    type: str
    text: str = ""
    placeholder_type: str | None = None
    width: int | None = None
    height: int | None = None
    char_limit_estimate: int | None = None
    rows: list[list[str]] = Field(default_factory=list)


class TemplateSlide(BaseModel):
    slide_index: int
    title: str | None = None
    elements: list[TemplateElement] = Field(default_factory=list)


class TemplateStructure(BaseModel):
    slide_width: int
    slide_height: int
    slide_count: int
    slides: list[TemplateSlide]


def analyze_template(path: str | Path) -> TemplateStructure:
    prs = Presentation(str(path))
    slides: list[TemplateSlide] = []
    for slide_index, slide in enumerate(prs.slides):
        elements: list[TemplateElement] = []

        def add_shape(shape, element_id: str) -> None:
            element_type = "shape"
            text = getattr(shape, "text", "") or ""
            placeholder_type = None
            if getattr(shape, "is_placeholder", False):
                element_type = "placeholder"
                placeholder_type = str(shape.placeholder_format.type)
            if getattr(shape, "has_table", False):
                element_type = "table"
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
            else:
                rows = []
            if getattr(shape, "has_chart", False):
                element_type = "chart"
            width = getattr(shape, "width", None)
            height = getattr(shape, "height", None)
            char_limit = None
            if width and height:
                # Rough estimate: assuming average character is 12x12 pixels/EMUs, etc. 
                # PowerPoint dimensions are in EMU (914400 EMU = 1 inch = 72 pt). 
                # Let's say a 12pt font is 1/6 inch high = 152400 EMU.
                # Average character width for 12pt is about 75000 EMU.
                char_limit = int((width / 80000) * (height / 180000))

            if text or rows or element_type in {"placeholder", "chart"}:
                elements.append(TemplateElement(id=element_id, type=element_type, text=text, placeholder_type=placeholder_type, rows=rows, width=width, height=height, char_limit_estimate=char_limit))
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for cell_index, cell in enumerate(row.cells):
                        elements.append(
                            TemplateElement(
                                id=f"{element_id}_cell_{row_index}_{cell_index}",
                                type="table_cell",
                                text=cell.text,
                            )
                        )
            if hasattr(shape, "shapes"):
                for child_index, child in enumerate(shape.shapes):
                    add_shape(child, f"{element_id}_group_{child_index}")

        for shape_index, shape in enumerate(slide.shapes):
            add_shape(shape, f"slide_{slide_index}_shape_{shape_index}")
        title = slide.shapes.title.text if slide.shapes.title else None
        slides.append(TemplateSlide(slide_index=slide_index, title=title, elements=elements))
    return TemplateStructure(slide_width=prs.slide_width, slide_height=prs.slide_height, slide_count=len(prs.slides), slides=slides)
